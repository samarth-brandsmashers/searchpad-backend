from fastapi import FastAPI, Response
from app.routes import video
from dotenv import load_dotenv
import os
from io import BytesIO
from google.cloud import texttospeech
from fastapi import FastAPI, Form, HTTPException , BackgroundTasks
from pptx import Presentation
import requests, time, json

import uvicorn
import openai
import cv2
import os
from gtts import gTTS
import tempfile
import numpy as np
import pyttsx3
from pydantic import BaseModel
from pptx.util import Inches, Pt
from fastapi.responses import FileResponse, StreamingResponse
import asyncio
from fastapi.middleware.cors import CORSMiddleware
from pptx.dml.color import RGBColor
from PIL import Image


# Load environment variables from .env file
load_dotenv()

app = FastAPI(
    title="AI Video Generation API",
    description="API for generating videos using AI based on user prompts.",
    version="1.0.0"
)

# Add CORS Middleware
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # Allow all origins; you can specify specific origins like ["http://example.com"]
    allow_credentials=True,
    allow_methods=["*"],  # Allow all methods; you can specify specific methods like ["GET", "POST"]
    allow_headers=["*"],  # Allow all headers; you can specify specific headers like ["X-Custom-Header"]
)

# Include video routes
app.include_router(video.router, prefix="/api/v1", tags=["Video Generation"])

@app.get("/")
def read_root():
    return {"message": "Welcome to the AI Video Generation API!"}

@app.post("/create-video-withai/")
def generate_video():
    try:
        print("Starting the video generation process...")  # Debugging message
        url = "https://modelslab.com/api/v6/video/text2video"
        payload = json.dumps({
            "key": "1ydXx9DDizhpyjvR9bHy5tzcaJlDxf1EzdAJ3EnZOSX8aVub2UZOzQ8EJMx3",
            "model_id": "zeroscope",
            "prompt": "Man walking on road",
            "negative_prompt": "low quality, blurry",
            "height": 320,
            "width": 512,
            "num_frames": 16,
            "num_inference_steps": 20,
            "guidance_scale": 7,
            "upscale_height": 640,
            "upscale_width": 1024,
            "upscale_strength": 0.6,
            "upscale_guidance_scale": 12,
            "upscale_num_inference_steps": 20,
            "output_type": "mp4",
            "webhook": None,
            "track_id": None
        })
        headers = {'Content-Type': 'application/json'}
        print("Sending request to the API...")  # Debugging message
        response = requests.post(url, headers=headers, data=payload)
        print(f"Response from API: {response.status_code}")  # Debugging message
        result = response.json()
        print(f"Result: {result}")  # Debugging message
        eta = result.get('eta', 60)  # Default to 60 seconds if not provided
        fetch_url = result.get('fetch_result', "")
        print(f"Waiting for {eta} seconds before fetching the result...")  # Debugging message
        time.sleep(eta)
        print(f"Fetching result from {fetch_url}...")  # Debugging message
        fetch_response = requests.post(fetch_url, headers=headers, data=json.dumps({}))
        if fetch_response.status_code == 200:
            print("Video processed successfully!")
            return {"message": "Video processed successfully!", "data": fetch_response.json(), "url": fetch_url}
        else:
            print(f"Failed to fetch the result: {fetch_response.status_code}")
            return {"error": "Failed to fetch the result", "status_code": fetch_response.status_code}
    except Exception as e:
        print(f"Error occurred: {str(e)}")
        return {"error": str(e)}

openai.api_key = '<openai_key>'

class TTSRequest(BaseModel):
    text: str
    voice_engine: str
    voice: str
    output_format: str
    emotion: str

@app.post("/text-to-speech/")
async def text_to_speech(tts_request: TTSRequest, background_tasks: BackgroundTasks):
    temp_file_path = None  # Initialize the temp_file_path variable

    try:
        USER_ID = "UI21VZi4vick3DNf2W2dcyu5Sdt1"
        API_KEY = "c3a3cc60b1384da2838862d2cd6f8b83"

        url = "https://api.play.ht/api/v2/tts/stream"
        headers = {
            'X-USER-ID': USER_ID,
            'AUTHORIZATION': API_KEY,
            'accept': 'audio/mpeg',
            'content-type': 'application/json'
        }

        data = {
            "text": tts_request.text,
            "voice_engine": tts_request.voice_engine,
            "voice": tts_request.voice,
            "output_format": tts_request.output_format,
            "emotion": tts_request.emotion
        }

        response = requests.post(url, headers=headers, json=data)

        if response.status_code != 200:
            raise HTTPException(status_code=response.status_code, detail=response.text)

        with tempfile.NamedTemporaryFile(delete=False, suffix='.mp3') as temp_file:
            temp_file.write(response.content)
            temp_file_path = temp_file.name  # Store the path of the created file

        # Schedule the deletion of the temporary file after the response is sent
        background_tasks.add_task(delete_file, temp_file_path)

        # Return the audio file as a FileResponse
        return FileResponse(temp_file_path, media_type='audio/mpeg', filename='result.mp3')

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


async def delete_file(file_path: str):
    await asyncio.sleep(5)  # Optional delay before deletion
    if os.path.exists(file_path):
        os.remove(file_path)

class ImageRequest(BaseModel):
    prompt: str

def sharpen_image(image_url: str) -> np.ndarray:
    # Step 1: Download the image from the provided URL
    response = requests.get(image_url)
    if response.status_code != 200:
        raise HTTPException(status_code=500, detail="Could not download image.")

    # Step 2: Convert the image content (bytes) into a NumPy array
    image_array = np.asarray(bytearray(response.content), dtype=np.uint8)

    # Step 3: Decode the image array to OpenCV format (BGR)
    image = cv2.imdecode(image_array, cv2.IMREAD_COLOR)
    if image is None:
        raise HTTPException(status_code=500, detail="Could not process image.")

    # Step 4: Apply the first sharpening using a high-pass filter
    kernel = np.array([[0, -1, 0],
                       [-1, 5, -1],
                       [0, -1, 0]])
    sharpened_image = cv2.filter2D(image, -1, kernel)

    # Step 5: Apply unsharp masking for more sharpening
    blurred = cv2.GaussianBlur(sharpened_image, (9, 9), 10)
    sharpened_final = cv2.addWeighted(sharpened_image, 1.5, blurred, -0.5, 0)

    # Step 6: (Optional) Apply adaptive sharpening using CLAHE
    # CLAHE (Contrast Limited Adaptive Histogram Equalization) enhances contrast in small areas
    lab = cv2.cvtColor(sharpened_final, cv2.COLOR_BGR2LAB)
    l_channel, a, b = cv2.split(lab)

    # Apply CLAHE to the L-channel (brightness)
    clahe = cv2.createCLAHE(clipLimit=3.0, tileGridSize=(8, 8))
    cl = clahe.apply(l_channel)

    # Merge the CLAHE enhanced L-channel back with A and B channels
    limg = cv2.merge((cl, a, b))
    final_image = cv2.cvtColor(limg, cv2.COLOR_LAB2BGR)

    return final_image


@app.post("/graphic/")
async def generate_image(request: ImageRequest):
    print(request)
    try:
        # Call OpenAI's DALL-E API to generate an image
        response = openai.Image.create(
            prompt=request.prompt,
            n=1,  # You can increase this to get multiple images
            size="1024x1024",  # Use the largest size available
            response_format="url"  # Ensure you get a URL for the generated image
        )
        
        image_url = response['data'][0]['url']

        # Sharpen the image
        # sharpened_image = sharpen_image(image_url)

        # Encode the sharpened image back to a format that can be sent as a response
        # _, buffer = cv2.imencode('.png', sharpened_image)
        # sharpened_image_bytes = buffer.tobytes()
        print(image_url)
        return Response(content=image_url, media_type="text/plain")
 
    except Exception as e:
        print(e)
        raise HTTPException(status_code=500, detail=str(e))
    
# Helper function to download image from URL
def download_image(image_url):
    response = requests.get(image_url)
    if response.status_code == 200:
        temp_file = tempfile.NamedTemporaryFile(delete=False, suffix=".png")
        with open(temp_file.name, 'wb') as f:
            f.write(response.content)
        return temp_file.name
    else:
        raise HTTPException(status_code=400, detail="Failed to download image.")

class Slide(BaseModel):
    prompt: str  

class PresentationRequest(BaseModel):
    slides: list[Slide]

@app.post("/generate-presentation/")
async def generate_presentation(prompt: str = Form(...), primary_color: str = Form(...), secondary_color: str = Form(...)):
    try:
        slide_count = 6
        print(f"Generating Presentation with {slide_count} slides")
        
        # Convert the user-provided color to RGB
        primary_rgb = hex_to_rgb(primary_color)
        secondary_rgb = hex_to_rgb(secondary_color)

        # Create a new PowerPoint presentation
        presentation = Presentation()

        for i in range(slide_count):
            # Generate a unique title, content, and image suggestion using OpenAI for each slide
            response = openai.ChatCompletion.create(
                model="gpt-3.5-turbo",
                messages=[{
                    "role": "user",
                    "content": f"Generate a slide for a PowerPoint presentation on the topic: {prompt}. "
                               f"For slide {i + 1}, generate the following format:\n"
                               "1. Title: (a unique title for the slide)\n"
                               "2. Content: (a brief description or explanation)\n"
                               "3. Image: (suggestion for an image related to the slide)"
                }]
            )
            ai_text = response['choices'][0]['message']['content']

            # Extract title, content, and image keywords with better error handling
            title_text = extract_text_by_label(ai_text, "Title:", default=f"Slide {i + 1} - Generated Title")
            content_text = extract_text_by_label(ai_text, "Content:", default="Content not generated.")
            image_keywords = extract_text_by_label(ai_text, "Image:", default="generic topic image")

            # Add a slide layout with a title and content layout
            slide_layout = presentation.slide_layouts[1]  # Title and Content layout
            slide = presentation.slides.add_slide(slide_layout)

            # Set slide background to primary color
            slide_bg = slide.background
            fill = slide_bg.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(*primary_rgb)

            # Set the AI-generated title
            title = slide.shapes.title
            title.text = title_text

            # Set the AI-generated content
            content = slide.placeholders[1]
            content.text = content_text

                    # Apply secondary color to content text
            for paragraph in content.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(16)  # Set a reasonable font size
                    run.font.color.rgb = RGBColor(*secondary_rgb)  # Apply secondary color to text
            content.text_frame.word_wrap = True  # Ensure text wraps within the text box    

                    # Apply secondary color to content text 
            for paragraph in content.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(16)  # Set a reasonable font size
                    run.font.color.rgb = RGBColor(*secondary_rgb)  # Apply secondary color to text
            content.text_frame.word_wrap = True  # Ensure text wraps within the text box    

            # Apply secondary color to content text
            for paragraph in content.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(16)  # Set a reasonable font size
                    run.font.color.rgb = RGBColor(*secondary_rgb)  # Apply secondary color to text
            content.text_frame.word_wrap = True  # Ensure text wraps within the text box

            # Fetch an image related to the keywords using an image API
            image_url = get_image_url_based_on_keywords(image_keywords)
            if image_url:
                try:
                    # Download the image
                    image_data = download_image(image_url)

                    # Open image with Pillow
                    image = Image.open(BytesIO(image_data))

                    # Resize image while maintaining the aspect ratio
                    max_width = 8  # Maximum width in inches
                    max_height = 2.5  # Maximum height in inches
                    image.thumbnail((max_width * 96, max_height * 96), Image.LANCZOS)  # 96 DPI for PowerPoint

                    # Save resized image to a BytesIO object
                    img_stream = BytesIO()
                    image.save(img_stream, format='PNG')
                    img_stream.seek(0)

                    # Add the resized image to the slide
                    left = Inches(2.5)  # Positioning the image from the left
                    top = Inches(4)  # Positioning the image below the content; adjust as necessary

                    # Add picture to slide
                    slide.shapes.add_picture(img_stream, left, top)

                except Exception as e:
                    print(f"Failed to download image for slide: {title_text} with keywords: {image_keywords}. Error: {str(e)}")
            else:
                print(f"No image found for keywords: {image_keywords}. Skipping image for slide: {title_text}.")

        # Save the presentation to a BytesIO object
        ppt_io = BytesIO()
        presentation.save(ppt_io)
        ppt_io.seek(0)

        # Return the PowerPoint file as a response
        return StreamingResponse(
            ppt_io,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": "attachment; filename=presentation.pptx"}
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

def extract_text_by_label(text: str, label: str, default: str = "") -> str:
    """
    Extracts the text following a specific label (e.g., 'Title:', 'Content:').
    If the label isn't found, it returns the default value.
    """
    try:
        if label in text:
            return text.split(label)[1].split("\n")[0].strip()
        else:
            return default
    except Exception as e:
        print(f"Error extracting {label}: {e}")
        return default

def hex_to_rgb(hex_color: str) -> tuple:
    """
    Converts a hex color string (e.g., "#RRGGBB") to an RGB tuple.
    """
    hex_color = hex_color.lstrip('#')
    return tuple(int(hex_color[i:i + 2], 16) for i in (0, 2, 4))

def get_image_url_based_on_keywords(keywords: str) -> str:
    url = f"https://api.pexels.com/v1/search?query={keywords}&per_page=1"
    headers = {"Authorization": '45V1GVw8okGICULQRW2K9npfogEmBcGixA8nTIUidOsxi85dzUhjmbPO'}  
    response = requests.get(url, headers=headers)
    if response.status_code == 200:
        data = response.json()
        if data['photos']:
            return data['photos'][0]['src']['medium']  # Get the medium-sized image URL
    return None

def download_image(image_url: str) -> bytes:
    print(f"Downloading image from URL: {image_url}")  # Added for debugging
    response = requests.get(image_url)
    if response.status_code == 200:
        return response.content
    else:
        raise Exception(f"Failed to download image from {image_url}")
